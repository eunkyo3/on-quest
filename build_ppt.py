# -*- coding: utf-8 -*-
"""On-Quest 최종발표자료 생성 (python-pptx)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

BODY = "맑은 고딕"
MONO = "Consolas"
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
ACCENT2 = RGBColor(0x2E, 0x75, 0xB6)
LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
DARK = RGBColor(0x26, 0x2A, 0x33)
GREY = RGBColor(0x7F, 0x7F, 0x7F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def kfont(run, name=BODY, size=18, bold=False, color=DARK, italic=False):
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    rpr = run._r.get_or_add_rPr()
    ea = rpr.find(qn('a:ea'))
    if ea is None:
        ea = rpr.makeelement(qn('a:ea'), {}); rpr.append(ea)
    ea.set('typeface', name)

def slide():
    return prs.slides.add_slide(BLANK)

def rect(s, l, t, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def textbox(s, l, t, w, h, lines, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of (text, size, bold, color, font) or list of such-> multiple runs per para if nested list."""
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        p.space_after = Pt(6)
        runs = ln if isinstance(ln, list) else [ln]
        for rspec in runs:
            text, size, bold, color, font = (rspec + (BODY,))[:5] if len(rspec) < 5 else rspec
            r = p.add_run(); r.text = text
            kfont(r, font, size, bold, color)
    return tb

def title_bar(s, title, num=None):
    rect(s, 0, 0, SW, Inches(1.12), ACCENT)
    rect(s, 0, Inches(1.12), SW, Inches(0.06), ACCENT2)
    textbox(s, Inches(0.55), Inches(0.16), Inches(11.5), Inches(0.85),
            [[(title, 28, True, WHITE)]], anchor=MSO_ANCHOR.MIDDLE)
    if num is not None:
        textbox(s, SW-Inches(1.2), Inches(0.30), Inches(0.9), Inches(0.5),
                [[(num, 14, True, WHITE)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def footer(s):
    textbox(s, Inches(0.55), SH-Inches(0.45), Inches(8), Inches(0.35),
            [[("On-Quest  ·  애니셀 정은교 (2025800216)", 10, False, GREY)]])

def bullets(s, items, l=Inches(0.7), t=Inches(1.5), w=Inches(12), h=Inches(5.4), size=18, gap=10):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(gap)
        lvl = 0
        if isinstance(it, tuple) and isinstance(it[0], int):
            lvl = it[0]; it = it[1]
        p.level = lvl
        mark = "•  " if lvl == 0 else "–  "
        if isinstance(it, tuple):
            lead, rest = it
            r = p.add_run(); r.text = mark + lead; kfont(r, BODY, size-(2*lvl), True, ACCENT if lvl==0 else ACCENT2)
            r2 = p.add_run(); r2.text = rest; kfont(r2, BODY, size-(2*lvl), False, DARK)
        else:
            r = p.add_run(); r.text = mark + it; kfont(r, BODY, size-(2*lvl), False, DARK)
    return tb

def table(s, headers, rows, l, t, w, h, fs=12, colw=None):
    gt = s.shapes.add_table(len(rows)+1, len(headers), l, t, w, h).table
    if colw:
        for i, cw in enumerate(colw): gt.columns[i].width = Inches(cw)
    for j, htext in enumerate(headers):
        c = gt.cell(0, j); c.fill.solid(); c.fill.fore_color.rgb = ACCENT
        c.margin_top = Pt(3); c.margin_bottom = Pt(3)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = htext; kfont(r, BODY, fs, True, WHITE)
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            c = gt.cell(i, j)
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
            c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if len(str(val)) < 14 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val); kfont(r, BODY, fs, False, DARK)
    return gt

def flowbox(s, l, t, w, h, title, sub, fill=ACCENT, tcolor=WHITE):
    sp = rect(s, l, t, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_top=Pt(2); tf.margin_bottom=Pt(2)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; kfont(r, BODY, 13, True, tcolor)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; kfont(r2, BODY, 9.5, False, tcolor)
    return sp

def arrow(s, l, t, w=Inches(0.5), label=None):
    a = rect(s, l, t, w, Inches(0.36), ACCENT2, shape=MSO_SHAPE.RIGHT_ARROW)
    if label:
        textbox(s, l-Inches(0.1), t-Inches(0.42), w+Inches(0.2), Inches(0.4),
                [[(label, 9, False, GREY)]], align=PP_ALIGN.CENTER)
    return a

# ============== Slide 1: 표지 ==============
s = slide()
rect(s, 0, 0, SW, SH, ACCENT)
rect(s, 0, Inches(4.55), SW, Inches(0.06), ACCENT2)
textbox(s, Inches(0.9), Inches(1.4), Inches(11.5), Inches(0.6),
        [[("2026-1학기 PBL 최종 발표", 18, True, RGBColor(0xBF,0xD3,0xEA))]])
textbox(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(1.6),
        [[("On-Quest", 60, True, WHITE)]])
textbox(s, Inches(0.95), Inches(3.55), Inches(11.5), Inches(0.8),
        [[("게임형(퀘스트형) 신입사원 온보딩 자동화 플랫폼", 22, True, RGBColor(0xDC,0xE7,0xF3))]])
textbox(s, Inches(0.95), Inches(4.9), Inches(11.5), Inches(2.0),
        [[("재직 기업 :  애니셀", 16, False, WHITE)],
         [("발표자 :  정은교 (학번 2025800216)  ·  개인 프로젝트", 16, False, WHITE)],
         [("기업현장교사 :  문용근        지도교수 :  김현우", 16, False, WHITE)],
         [("2026년 6월", 14, False, RGBColor(0xBF,0xD3,0xEA))]])

# ============== Slide 2: 목차 ==============
s = slide(); title_bar(s, "목차", "01")
bullets(s, [
 ("01. 문제 정의 & 필요성", ""),
 ("02. 프로젝트 소개 — On-Quest", ""),
 ("03. 주안점 & 직무 연관성", ""),
 ("04. 작동 흐름 & 주요 기능", ""),
 ("05. 시스템 아키텍처 & 워크플로우", ""),
 ("06. 기술 스택 & 핵심 구현(보안)", ""),
 ("07. 검증 (기능·성능)", ""),
 ("08. 시연 영상", ""),
 ("09. 기대 효과 & 마무리", ""),
], t=Inches(1.55), size=20, gap=12)
footer(s)

# ============== Slide 3: 문제 정의 & 필요성 ==============
s = slide(); title_bar(s, "문제 정의 & 필요성", "02")
rect(s, Inches(0.7), Inches(1.6), Inches(5.7), Inches(4.6), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, Inches(1.0), Inches(1.85), Inches(5.1), Inches(4.2),
        [[("문제", 22, True, ACCENT)],
         [("", 8, False, DARK)],
         [("• 신입사원이 입사 후 ", 17, False, DARK), ("무엇을 해야 할지 몰라", 17, True, DARK), (" 헤맨다", 17, False, DARK)],
         [("• 사수(선임)는 ", 17, False, DARK), ("자기 업무로 바빠", 17, True, DARK), (" 신경 쓰기 어렵다", 17, False, DARK)],
         [("• 결과적으로 신입은 업무가 주어지기 전까지 ", 17, False, DARK), ("대기·방치", 17, True, DARK)],
         [("• 조직 적응이 지연됨", 17, False, DARK)]])
rect(s, Inches(6.9), Inches(1.6), Inches(5.7), Inches(4.6), ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, Inches(7.2), Inches(1.85), Inches(5.1), Inches(4.2),
        [[("필요성", 22, True, WHITE)],
         [("", 8, False, WHITE)],
         [("• 입사 ", 17, False, WHITE), ("즉시 해야 할 일을 스스로 확인", 17, True, WHITE)],
         [("• 부여–수행–검토–피드백을 ", 17, False, WHITE), ("자동화된 사이클", 17, True, WHITE), ("로", 17, False, WHITE)],
         [("• 신입의 ", 17, False, WHITE), ("빠른 조직 적응", 17, True, WHITE), (" 지원", 17, False, WHITE)],
         [("• 관리자의 ", 17, False, WHITE), ("온보딩 부담 경감", 17, True, WHITE)]])
footer(s)

# ============== Slide 4: 프로젝트 소개 ==============
s = slide(); title_bar(s, "프로젝트 소개 — On-Quest", "03")
textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(1.2),
        [[("\"신입사원의 온보딩을, 게임 퀘스트처럼.\"", 24, True, ACCENT)],
         [("회사가 신입에게 줄 할 일을 '퀘스트'로 정의하고 부여→수행→인증샷 제출→검토→피드백/완료를 하나의 워크플로우로 자동화", 15, False, DARK)]])
bullets(s, [
 ("게임화(Gamification) : ", "할 일을 상태(Status)를 가진 '퀘스트'로 모델링 — 단순 체크리스트가 아닌 워크플로우"),
 ("실시간 알림 : ", "각 단계 이벤트를 n8n→Slack으로 자동 통지 — 별도 채근 없이 사이클 진행"),
 ("3단계 권한 : ", "슈퍼관리자 / 관리자 / 신입사원, 회사코드 기준 데이터 완전 분리(멀티테넌시)"),
 ("회사 첫 가입자 = 슈퍼관리자 : ", "초기 세팅 없이 바로 구성원·권한 관리 시작"),
], t=Inches(3.0), size=17, gap=12)
footer(s)

# ============== Slide 5: 주안점 & 직무 연관성 ==============
s = slide(); title_bar(s, "주안점 & 직무 연관성", "04")
textbox(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.5),
        [[("주안점 (Key Point)", 19, True, ACCENT)]])
kp = [("게임화", "상태형 퀘스트 모델"), ("시스템 통합", "서버↔n8n↔Slack"),
      ("보안·신뢰성", "HMAC·JWT·감사로그"), ("멀티테넌시", "회사별 데이터 격리")]
x = Inches(0.7)
for i,(tt,ss) in enumerate(kp):
    flowbox(s, x, Inches(2.0), Inches(2.85), Inches(1.15), tt, ss, fill=ACCENT2 if i%2 else ACCENT)
    x += Inches(3.0)
rect(s, Inches(0.7), Inches(3.55), Inches(11.93), Inches(0.04), ACCENT2)
textbox(s, Inches(0.7), Inches(3.75), Inches(12), Inches(0.5),
        [[("직무 연관성 — 시스템 통합 역량", 19, True, ACCENT)]])
bullets(s, [
 ("서로 다른 시스템 연결 : ", "백엔드(NestJS) ↔ 자동화(n8n) ↔ 메신저(Slack)를 이벤트 기반으로 유기적 연동"),
 ("서드파티 API 연동 + 인프라 자동화 : ", "기업 업무 효율화를 위한 백엔드 개발자의 핵심 '시스템 통합 역량'과 직결"),
 ("단순 호출이 아닌 보안 연동 : ", "모든 웹훅 페이로드를 HMAC-SHA256으로 서명·검증"),
], t=Inches(4.3), size=16, gap=10)
footer(s)

# ============== Slide 6: 작동 흐름 ==============
s = slide(); title_bar(s, "작동 흐름 (역할별)", "05")
rect(s, Inches(0.7), Inches(1.55), Inches(5.7), Inches(4.7), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, Inches(1.0), Inches(1.8), Inches(5.1), Inches(0.6),[[("관리자", 20, True, ACCENT)]])
bullets(s, [("퀘스트 발행 (단건 / CSV 일괄)",""),
            ("제출된 인증샷 검토",""),
            ("승인(완료) 또는 피드백과 함께 반려",""),
            ("거부된 퀘스트 재개봉·재배정","")],
        l=Inches(1.0), t=Inches(2.45), w=Inches(5.1), h=Inches(3.6), size=15, gap=12)
rect(s, Inches(6.9), Inches(1.55), Inches(5.7), Inches(4.7), ACCENT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
textbox(s, Inches(7.2), Inches(1.8), Inches(5.1), Inches(0.6),[[("신입사원", 20, True, WHITE)]])
tb = s.shapes.add_textbox(Inches(7.2), Inches(2.45), Inches(5.1), Inches(3.6)); tf=tb.text_frame; tf.word_wrap=True
for i,txt in enumerate(["할당된 퀘스트 확인","착수 (수행 거부 시 사유 입력)","과제 완료 후 인증샷 제출","반려 시 피드백 보고 → 보완 재제출"]):
    p = tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(12)
    r=p.add_run(); r.text="•  "+txt; kfont(r, BODY, 15, False, WHITE)
footer(s)

# ============== Slide 7: 주요 기능 ==============
s = slide(); title_bar(s, "주요 기능 (기능 요구사항)", "06")
table(s, ["역할","기능","설명"],
 [["공통","회원가입/로그인","회사코드 기반 가입·JWT 인증, 첫 가입자 슈퍼관리자 자동 지정"],
  ["슈퍼관리자","구성원·권한 관리","역할 변경, 소유권 이양, 감사 로그 조회"],
  ["관리자","퀘스트 발행","단건 + CSV 일괄(최대 100건), 마감·담당자 지정"],
  ["관리자","검토","승인(완료)·피드백 반려·거부 건 재개봉"],
  ["신입사원","퀘스트 수행","확인·착수·거부(사유)·인증샷 제출/재제출"],
  ["공통","증빙 파일","이미지·PDF 업로드, 미리보기, Slack 서명 링크 공유"],
  ["시스템","실시간 알림","발행·제출·검토·마감 임박/초과 → Slack"],
  ["관리자","진행 현황","진행률 대시보드·담당자별 통계·CSV 내보내기"]],
 Inches(0.7), Inches(1.55), Inches(11.93), Inches(5.0), fs=13, colw=[1.7,2.3,7.93])
footer(s)

# ============== Slide 8: 아키텍처 ==============
s = slide(); title_bar(s, "시스템 아키텍처", "07")
textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
        [[("전 구성요소를 Docker Compose로 기동 · Nginx가 /api를 백엔드로 프록시", 14, False, GREY)]])
y = Inches(2.5)
flowbox(s, Inches(0.55), y, Inches(2.1), Inches(1.2), "사용자", "브라우저 · React", fill=RGBColor(0x55,0x55,0x55))
arrow(s, Inches(2.75), Inches(2.9), Inches(0.5), "HTTPS")
flowbox(s, Inches(3.35), y, Inches(2.1), Inches(1.2), "Frontend", "React+Vite (Nginx)", fill=ACCENT2)
arrow(s, Inches(5.55), Inches(2.9), Inches(0.5), "/api")
flowbox(s, Inches(6.15), y, Inches(2.3), Inches(1.2), "Backend (NestJS)", "auth·quest·audit·automation", fill=ACCENT)
flowbox(s, Inches(6.15), Inches(4.35), Inches(2.3), Inches(1.0), "PostgreSQL 16", "Prisma ORM", fill=RGBColor(0x33,0x67,0x9A))
# arrow backend->db (down)
da = rect(s, Inches(7.05), Inches(3.7), Inches(0.5), Inches(0.65), ACCENT2, shape=MSO_SHAPE.DOWN_ARROW)
arrow(s, Inches(8.55), Inches(2.9), Inches(0.55), "HMAC")
flowbox(s, Inches(9.2), y, Inches(1.7), Inches(1.2), "n8n", "워크플로우 자동화", fill=RGBColor(0xC0,0x5A,0x2A))
arrow(s, Inches(10.95), Inches(2.9), Inches(0.5), "알림")
flowbox(s, Inches(11.5), y, Inches(1.3), Inches(1.2), "Slack", "실시간 알림", fill=RGBColor(0x3F,0x0F,0x3F))
textbox(s, Inches(6.15), Inches(5.45), Inches(6.6), Inches(0.6),
        [[("Backend → DB(Prisma) 저장,  Backend → n8n(HMAC 서명 webhook) → Slack 알림", 12, False, GREY)]])
footer(s)

# ============== Slide 9: 퀘스트 워크플로우 ==============
s = slide(); title_bar(s, "퀘스트 상태 워크플로우", "08")
states = [("대기","0",ACCENT),("착수","1",ACCENT2),("검토대기","2",ACCENT2),("완료","3",RGBColor(0x2E,0x7D,0x32))]
x = Inches(0.6)
for i,(nm,code,col) in enumerate(states):
    flowbox(s, x, Inches(2.2), Inches(2.3), Inches(1.0), nm, "status "+code, fill=col)
    if i < 3: arrow(s, x+Inches(2.35), Inches(2.52), Inches(0.45),
                    ["발행/착수","제출","승인"][i])
    x += Inches(2.95)
flowbox(s, Inches(3.55), Inches(4.2), Inches(2.3), Inches(0.95), "반려", "status 4", fill=RGBColor(0xC6,0x2D,0x2D))
flowbox(s, Inches(0.6), Inches(4.2), Inches(2.3), Inches(0.95), "거부됨", "status 5", fill=RGBColor(0x8A,0x8A,0x8A))
textbox(s, Inches(0.6), Inches(5.45), Inches(12), Inches(1.4),
        [[("• 검토대기 → 반려(피드백) → 보완 후 재제출 → 다시 검토대기", 15, False, DARK)],
         [("• 대기/착수 → (사원 거부) → 거부됨 → (관리자 재개봉) → 대기", 15, False, DARK)],
         [("• 모든 전이는 '기대 상태일 때만 갱신'되는 원자적 연산 (동시 요청 시 409 Conflict)", 14, True, ACCENT)]])
footer(s)

# ============== Slide 10: 기술 스택 ==============
s = slide(); title_bar(s, "기술 스택 & 계획 대비 변경", "09")
table(s, ["구분","채택 기술"],
 [["언어","TypeScript (프론트·백 공통)"],
  ["프론트엔드","React 18(Vite) · Zustand · React Router · Axios"],
  ["백엔드","NestJS 10 · Passport-JWT · class-validator · schedule/throttler"],
  ["DB / ORM","PostgreSQL 16 / Prisma"],
  ["자동화·연동","n8n · Slack API · HMAC-SHA256"],
  ["인프라","Docker Compose · Nginx · GitHub Actions · Jest"]],
 Inches(0.6), Inches(1.55), Inches(7.2), Inches(4.4), fs=12.5, colw=[1.5,5.7])
textbox(s, Inches(8.1), Inches(1.55), Inches(4.7), Inches(0.5),[[("계획 → 실제 (의사결정)", 15, True, ACCENT)]])
table(s, ["계획","실제"],
 [["TypeORM","Prisma"],
  ["TanStack Query","Zustand"],
  ["Tailwind","순수 CSS"],
  ["AWS S3","DB BLOB 저장"]],
 Inches(8.1), Inches(2.1), Inches(4.7), Inches(2.5), fs=12.5, colw=[2.35,2.35])
textbox(s, Inches(8.1), Inches(4.8), Inches(4.7), Inches(1.4),
        [[("→ 외부 의존성·비용 최소화,", 12, False, GREY)],
         [("   타입 안정성·단순성 우선", 12, False, GREY)]])
footer(s)

# ============== Slide 11: 핵심 구현·보안 ==============
s = slide(); title_bar(s, "핵심 구현 — 보안 & 연동", "10")
cards = [
 ("웹훅 HMAC 서명/검증", "페이로드를 정규 직렬화 후 HMAC-SHA256 서명 → n8n이 재계산 검증. 위·변조/비인가 호출 차단 (±5분 타임스탬프)"),
 ("JWT · 토큰 폐기", "access(1h)+refresh(7d). refresh는 HttpOnly 쿠키, 로그아웃 시 tokenVersion 증가로 즉시 무효화"),
 ("원자적 상태 전이", "'기대 상태인 행만' 조건부 갱신으로 선점, 실패 시 409 → 동시성·중복 알림 방지"),
 ("증빙 보안", "허용 형식만 업로드(MIME 화이트리스트)+nosniff, 외부 공유는 HMAC 서명 토큰(timingSafeEqual)"),
 ("멀티테넌시·RBAC", "회사코드 스코핑으로 테넌트 격리, 역할 기반 접근제어, 인증 레이트리밋"),
 ("감사 로그·CSV 보안", "주요 행위 기록, CSV 내보내기 수식 인젝션(=,+,-,@) 무력화"),
]
positions = [(Inches(0.6),Inches(1.55)),(Inches(4.62),Inches(1.55)),(Inches(8.64),Inches(1.55)),
             (Inches(0.6),Inches(4.0)),(Inches(4.62),Inches(4.0)),(Inches(8.64),Inches(4.0))]
for (cx,cy),(ct,cd) in zip(positions, cards):
    rect(s, cx, cy, Inches(3.8), Inches(2.25), LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    textbox(s, cx+Inches(0.2), cy+Inches(0.15), Inches(3.45), Inches(2.0),
            [[(ct, 14.5, True, ACCENT)], [("", 4, False, DARK)], [(cd, 11.5, False, DARK)]])
footer(s)

# ============== Slide 12: 검증 ==============
s = slide(); title_bar(s, "검증 (기능·성능 평가)", "11")
textbox(s, Inches(0.6), Inches(1.4), Inches(6), Inches(0.45),[[("기능 동작 체크리스트 (발췌)", 15, True, ACCENT)]])
table(s, ["검증 항목","적/부"],
 [["회사 첫 가입자 슈퍼관리자 지정","적"],
  ["단건/CSV 일괄 발행 + 알림","적"],
  ["착수 → 증빙 제출 → 검토","적"],
  ["반려(피드백) → 재제출","적"],
  ["마감 임박/초과 자동 알림","적"],
  ["동시 검토 경합 → 409 처리","적"],
  ["권한 외 접근 차단 / 서명 검증","적"]],
 Inches(0.6), Inches(1.9), Inches(6.4), Inches(4.2), fs=12.5, colw=[5.2,1.2])
textbox(s, Inches(7.3), Inches(1.4), Inches(5.3), Inches(0.45),[[("품질·성능 지표", 15, True, ACCENT)]])
bullets(s, [
 ("테스트 : ", "6 스위트 / 36 테스트 전부 통과"),
 ("CI : ", "타입체크·린트·테스트·빌드 green"),
 ("알림 지연(목표) : ", "이벤트 후 3초 이내 (비동기·2.5s 타임아웃)"),
 ("정적 분석 : ", "TypeScript 타입체크·ESLint 통과"),
 ("가용성 : ", "헬스체크 + DB 프로브 타임아웃(3s)"),
], l=Inches(7.3), t=Inches(1.95), w=Inches(5.3), h=Inches(4), size=14.5, gap=14)
footer(s)

# ============== Slide 13: 시연 영상 전환 ==============
s = slide()
rect(s, 0, 0, SW, SH, ACCENT)
textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
        [[("🎬  시연 영상", 44, True, WHITE)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.0), Inches(11.3), Inches(1.2),
        [[("회원가입 → 퀘스트 발행 → Slack 알림 → 착수·제출 → 검토(반려/승인)", 18, False, RGBColor(0xDC,0xE7,0xF3))],
         [("→ 마감 알림 → 진행률 대시보드", 18, False, RGBColor(0xDC,0xE7,0xF3))]], align=PP_ALIGN.CENTER)

# ============== Slide 14: 기대 효과 & 마무리 ==============
s = slide(); title_bar(s, "기대 효과 & 향후 계획", "12")
textbox(s, Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),[[("기대 효과", 19, True, ACCENT)]])
bullets(s, [
 ("신입 : ", "사수가 일을 주기 전 공백 없이 즉시 할 일 시작 → 빠른 적응"),
 ("관리자 : ", "진행 현황을 대시보드로 한눈에, 반복적 전달·확인 커뮤니케이션 감소"),
 ("조직 : ", "온보딩 표준화 + 알림 자동화로 운영 효율 향상"),
], t=Inches(2.0), size=17, gap=10)
textbox(s, Inches(0.7), Inches(4.2), Inches(12), Inches(0.5),[[("향후 계획", 19, True, ACCENT)]])
bullets(s, [
 ("증빙 저장소를 외부 오브젝트 스토리지(S3 등)로 분리 (트래픽 확장 대비)",""),
 ("알림 채널 확장 (이메일·문자), 배지·레벨 등 게임화 요소 강화",""),
], t=Inches(4.7), size=16, gap=10)
footer(s)

# ============== Slide 15: 감사합니다 ==============
s = slide()
rect(s, 0, 0, SW, SH, ACCENT)
rect(s, 0, Inches(4.4), SW, Inches(0.05), ACCENT2)
textbox(s, Inches(1), Inches(2.6), Inches(11.3), Inches(1.2),
        [[("감사합니다", 48, True, WHITE)]], align=PP_ALIGN.CENTER)
textbox(s, Inches(1), Inches(4.7), Inches(11.3), Inches(1.5),
        [[("On-Quest  ·  게임형 신입사원 온보딩 자동화 플랫폼", 18, False, RGBColor(0xDC,0xE7,0xF3))],
         [("애니셀  정은교 (2025800216)", 16, False, WHITE)]], align=PP_ALIGN.CENTER)

prs.save(r"D:\최종발표자료.pptx")
print("SAVED", len(prs.slides.__iter__.__self__._sldIdLst))
